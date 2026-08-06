"""Build the detailed MetaGatherer deck and captioned demo slideshow."""

from __future__ import annotations

import subprocess
from pathlib import Path

from build_presentation import (
    ASSETS,
    BLUE,
    CYAN,
    GOLD,
    GREEN,
    MUTED,
    PANEL,
    PANEL_2,
    RED,
    WHITE,
    H,
    W,
    add_bg,
    add_metric,
    box,
    font,
    footer,
    title,
    txt,
)
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
PPTX_OUT = ROOT / "MetaGatherer_detailed_ru.pptx"
VIDEO_OUT = ROOT / "MetaGatherer_demo_ru.mp4"
VIDEO_DIR = ASSETS / "video"
FLOW_IMAGE = ASSETS / "architecture-flow.png"
AVATAR_IMAGE = ASSETS / "bot-avatar.png"
CARDS_IMAGE = ASSETS / "mtg-card-pair.png"


def add_picture_contain(slide, path: Path, x, y, w, h, bg="E9ECF1"):
    box(slide, x, y, w, h, bg, line="323A47")
    with Image.open(path) as im:
        ratio = im.width / im.height
    area = w / h
    if ratio > area:
        pw, ph = w - 0.12, (w - 0.12) / ratio
        px, py = x + 0.06, y + (h - ph) / 2
    else:
        ph, pw = h - 0.12, (h - 0.12) * ratio
        px, py = x + (w - pw) / 2, y + 0.06
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))


def build_flow_image():
    """High-level product architecture: one tournament from creation to publication."""
    im = Image.new("RGB", (1920, 1080), "#10141C")
    d = ImageDraw.Draw(im)
    d.text((100, 64), "КАК ПРОХОДИТ ТУРНИР", font=font(25, True), fill="#D6A928")
    d.text((100, 112), "От создания до результатов в клубном чате", font=font(54, True), fill="#F7F3E8")
    d.text(
        (100, 186),
        "Сквозной сценарий: от людей и турнирного сервиса — к общей истории меты",
        font=font(25),
        fill="#AAB3C2",
    )

    nodes = [
        (110, 315, "1", "Бот создаёт\nтурнир", "дата · клуб · формат", "#5591E6"),
        (665, 315, "2", "Уведомляет\nигроков", "opt-in и безопасная доставка", "#4FC3B3"),
        (1220, 315, "3", "Игроки\nзаписываются", "участник → колода", "#D6A928"),
        (1220, 690, "4", "Организатор\nдополняет", "закрывает пропуски", "#E4553D"),
        (665, 690, "5", "AetherHub отдаёт\nрезультаты", "пары · очки · места", "#55C787"),
        (110, 690, "6", "Бот публикует\nфинал", "мета и стендинги в чатах клубов", "#D6A928"),
    ]
    for x, y, n, head, body, color in nodes:
        d.rounded_rectangle((x, y, x + 430, y + 225), 30, fill="#202837", outline=color, width=5)
        d.ellipse((x + 24, y + 24, x + 78, y + 78), fill=color)
        nb = d.textbbox((0, 0), n, font=font(24, True))
        d.text((x + 51 - (nb[2] - nb[0]) / 2, y + 34), n, font=font(24, True), fill="#10141C")
        hy = y + 30
        for line in head.split("\n"):
            d.text((x + 102, hy), line, font=font(29, True), fill="#F7F3E8")
            hy += 36
        d.text((x + 28, y + 174), body, font=font(21), fill="#AAB3C2")

    # Snake-shaped process arrows.
    for x1, x2, y in [(540, 665, 427), (1095, 1220, 427)]:
        d.line((x1, y, x2 - 20, y), fill="#758092", width=8)
        d.polygon([(x2 - 22, y - 14), (x2, y), (x2 - 22, y + 14)], fill="#758092")
    for x1, x2, y in [(1220, 1095, 802), (665, 540, 802)]:
        d.line((x1, y, x2 + 20, y), fill="#758092", width=8)
        d.polygon([(x2 + 22, y - 14), (x2, y), (x2 + 22, y + 14)], fill="#758092")
    d.line((1435, 540, 1435, 670), fill="#758092", width=8)
    d.polygon([(1421, 650), (1435, 675), (1449, 650)], fill="#758092")
    # The published tournament also becomes part of cumulative statistics.
    d.line((325, 915, 325, 991), fill="#758092", width=8)
    d.line((325, 991, 450, 991), fill="#758092", width=8)
    d.polygon([(448, 977), (472, 991), (448, 1005)], fill="#758092")
    d.rounded_rectangle((470, 944, 1810, 1038), 28, fill="#171D28", outline="#4FC3B3", width=4)
    d.ellipse((500, 963, 558, 1021), fill="#4FC3B3")
    d.text((519, 974), "7", font=font(23, True), fill="#10141C")
    d.text((590, 962), "Турнир попадает в общую статистику", font=font(27, True), fill="#F7F3E8")
    d.text((590, 1002), "популярность колод · доля меты · матчи · винрейты", font=font(19), fill="#AAB3C2")
    im.save(FLOW_IMAGE, quality=95)


def build_supporting_images():
    """Crop the real bot avatar and compose two real cards for slides/video."""
    with Image.open(ASSETS / "real-telegram-recap.png").convert("RGBA") as source:
        avatar = source.crop((24, 8, 104, 88)).resize((420, 420))
        mask = Image.new("L", avatar.size, 0)
        ImageDraw.Draw(mask).ellipse((4, 4, 416, 416), fill=255)
        out = Image.new("RGBA", avatar.size, (0, 0, 0, 0))
        out.paste(avatar, (0, 0), mask)
        out.save(AVATAR_IMAGE)

    canvas = Image.new("RGB", (1500, 900), "#10141C")
    for path, x, angle in [
        (ASSETS / "card-tolarian-terror.jpg", 285, -5),
        (ASSETS / "card-experimental-synthesizer.jpg", 795, 5),
    ]:
        with Image.open(path).convert("RGB") as card:
            card.thumbnail((500, 740))
            card = card.rotate(angle, expand=True, resample=Image.Resampling.BICUBIC)
            canvas.paste(card, (x, 78), card if card.mode == "RGBA" else None)
    canvas.save(CARDS_IMAGE, quality=95)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    s.shapes.add_picture(str(AVATAR_IMAGE), Inches(0.72), Inches(0.34), width=Inches(0.62), height=Inches(0.62))
    txt(s, "METAGATHERER", 1.48, 0.52, 4, 0.3, 10, GOLD, True)
    txt(s, "Инфраструктура локального\nтурнира в Telegram", 0.72, 1.12, 7.4, 1.55, 34, WHITE, True)
    txt(
        s,
        "Как убрать ручную координацию, собрать качественные данные\nи вернуть их игрокам в полезном виде.",
        0.76,
        2.92,
        6.7,
        0.9,
        16,
        MUTED,
    )
    add_metric(s, 0.76, 4.35, "1 чат", "единая точка входа", CYAN, 1.8)
    add_metric(s, 2.75, 4.35, "4 стадии", "жизненный цикл турнира", GOLD, 1.95)
    add_metric(s, 4.90, 4.35, "3 источника", "данные и результаты", BLUE, 2.1)
    txt(
        s, "Кейс: Pauper — доступный соревновательный формат Magic: The Gathering", 0.8, 6.0, 6.5, 0.55, 15, WHITE, True
    )
    s.shapes.add_picture(str(ASSETS / "real-telegram-recap.png"), Inches(8.55), Inches(0.62), height=Inches(6.56))
    footer(s, 1)

    # 2 audience context
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Контекст без MTG", "Что происходит на обычном клубном турнире")
    steps = [
        ("До", "Люди решают, придут ли,\nи записываются", BLUE),
        ("Старт", "Каждый сообщает,\nкакой колодой играет", CYAN),
        ("Раунды", "Система формирует пары,\nигроки находят соперников", GOLD),
        ("Финиш", "Появляются места,\nочки и статистика", GREEN),
    ]
    for i, (head, body, c) in enumerate(steps):
        x = 0.7 + i * 3.12
        box(s, x, 2.25, 2.8, 2.3, PANEL, line=c)
        txt(s, str(i + 1), x + 0.2, 2.48, 0.45, 0.3, 13, c, True)
        txt(s, head, x + 0.2, 2.96, 2.35, 0.35, 20, WHITE, True)
        txt(s, body, x + 0.2, 3.54, 2.3, 0.65, 13, MUTED)
    txt(
        s,
        "Колода здесь — аналог выбранной стратегии/конфигурации. Метагейм — распределение этих стратегий в сообществе.",
        1.0,
        5.32,
        11.3,
        0.7,
        18,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    txt(
        s,
        "Даже маленькое событие создаёт много связанного состояния и ошибок синхронизации.",
        1.2,
        6.25,
        10.9,
        0.35,
        13,
        GOLD,
        align=PP_ALIGN.CENTER,
    )
    s.shapes.add_picture(str(ASSETS / "card-tolarian-terror.jpg"), Inches(0.68), Inches(4.73), height=Inches(1.52))
    s.shapes.add_picture(
        str(ASSETS / "card-experimental-synthesizer.jpg"), Inches(11.54), Inches(4.73), height=Inches(1.52)
    )
    footer(s, 2)

    # 3 pain
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Проблема", "До автоматизации данные рассыпаны по людям и сервисам")
    pains = [
        ("Telegram", "кто придёт? кто записался?", BLUE),
        ("Организатор", "у кого какая колода?", RED),
        ("AetherHub", "пары и финальные места", GOLD),
        ("Таблицы", "экспорт и история", GREEN),
    ]
    for i, (h, b, c) in enumerate(pains):
        x = 0.75 + i * 3.1
        box(s, x, 2.15, 2.75, 1.35, PANEL, line=c)
        txt(s, h, x + 0.18, 2.42, 2.38, 0.3, 17, WHITE, True)
        txt(s, b, x + 0.18, 2.92, 2.38, 0.25, 11, MUTED)
        if i < 3:
            txt(s, "↛", x + 2.78, 2.6, 0.3, 0.3, 18, RED, True, align=PP_ALIGN.CENTER)
    box(s, 1.55, 4.35, 10.2, 1.55, PANEL_2)
    txt(s, "Цена ручного процесса", 1.85, 4.62, 2.6, 0.3, 13, GOLD, True)
    txt(
        s,
        "пропущенные записи · дубли имён · забытые пары · несведённые результаты · нулевая история",
        1.85,
        5.12,
        9.55,
        0.4,
        18,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 3)

    # 4 player journey
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Игрок", "Запись занимает несколько нажатий")
    add_picture_contain(s, ASSETS / "real-registration.jpg", 0.7, 1.95, 5.8, 4.55)
    add_picture_contain(s, ASSETS / "real-choose-deck.jpg", 6.82, 1.95, 5.8, 4.55)
    s.shapes.add_picture(str(ASSETS / "card-tolarian-terror.jpg"), Inches(11.34), Inches(4.84), height=Inches(1.42))
    txt(s, "1 · выбрать активный турнир", 1.0, 6.63, 5.15, 0.26, 12, CYAN, True, align=PP_ALIGN.CENTER)
    txt(s, "2 · выбрать знакомую колоду или свой вариант", 7.05, 6.63, 5.3, 0.26, 12, GOLD, True, align=PP_ALIGN.CENTER)
    footer(s, 4, "Реальные кадры Telegram Desktop")

    # 5 admin journey
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Организатор", "Бот показывает незаполненные места и даёт закрыть их на ходу")
    add_picture_contain(s, ASSETS / "real-admin-fill.jpg", 0.75, 1.9, 6.4, 4.95)
    bullets = [
        ("Один экран состояния", "кто зарегистрирован и у кого нет колоды"),
        ("Коллаборативный сбор", "колоду может записать игрок, админ или оппонент"),
        ("Безопасные действия", "подтверждения для удаления, раскрытия и закрытия"),
    ]
    for i, (h, b) in enumerate(bullets):
        y = 2.05 + i * 1.42
        box(s, 7.55, y, 4.95, 1.1, PANEL)
        txt(s, h, 7.8, y + 0.18, 4.45, 0.28, 15, GOLD, True)
        txt(s, b, 7.8, y + 0.57, 4.42, 0.32, 11, MUTED)
    footer(s, 5, "Реальный debug-сценарий; уведомления не отправляются другим пользователям")

    # 6 state machine
    s = prs.slides.add_slide(blank)
    add_bg(s)
    s.shapes.add_picture(str(FLOW_IMAGE), 0, 0, width=Inches(W), height=Inches(H))

    # 7 state machine
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Модель", "Турнир — явный конечный автомат, а не набор команд")
    states = [
        ("REGISTRATION", "запись и колоды", BLUE),
        ("ONGOING", "раунды и пары", CYAN),
        ("CLOSED", "итоги и экспорт", GREEN),
    ]
    for i, (h, b, c) in enumerate(states):
        x = 1.35 + i * 4.0
        box(s, x, 2.35, 3.2, 1.52, PANEL, line=c)
        txt(s, h, x + 0.12, 2.7, 2.96, 0.28, 13, c, True, align=PP_ALIGN.CENTER)
        txt(s, b, x + 0.12, 3.22, 2.96, 0.25, 11, WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            txt(s, "→", x + 3.35, 2.88, 0.28, 0.3, 18, "738094", True)
    rules = [
        ("Инвариант", "один активный турнир на чат"),
        ("Гейты", "действия разрешены только в нужной стадии"),
        ("Идемпотентность", "повторный импорт не создаёт дубли"),
        ("Best effort", "внешняя ошибка не ломает закрытие"),
    ]
    for i, (h, b) in enumerate(rules):
        x = 0.82 + (i % 2) * 6.05
        y = 4.55 + (i // 2) * 0.92
        txt(s, h, x, y, 1.3, 0.25, 11, GOLD, True)
        txt(s, b, x + 1.45, y, 4.2, 0.28, 13, WHITE)
    footer(s, 7, "TournamentService + typed domain errors")

    # 8 architecture
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Архитектура", "Чистое ядро отделено от Telegram и интеграций")
    layers = [
        ("Telegram wrappers", "Update → primitives", BLUE),
        ("Pure handlers", "HandlerResult", CYAN),
        ("Services", "rules + orchestration", GOLD),
        ("SQLAlchemy", "models + migrations", GREEN),
    ]
    for i, (h, b, c) in enumerate(layers):
        y = 1.95 + i * 1.14
        box(s, 0.8, y, 5.0, 0.86, PANEL, line=c)
        txt(s, h, 1.05, y + 0.19, 2.65, 0.25, 15, WHITE, True)
        txt(s, b, 3.55, y + 0.21, 1.95, 0.22, 11, MUTED, align=PP_ALIGN.RIGHT)
        if i < 3:
            txt(s, "↓", 3.0, y + 0.88, 0.35, 0.22, 15, "758092", True, align=PP_ALIGN.CENTER)
    box(s, 6.4, 1.95, 5.95, 4.22, PANEL_2)
    txt(s, "Внешний контур", 6.72, 2.23, 2.0, 0.3, 16, GOLD, True)
    integrations = [
        ("AetherHub", "паринги и стендинги"),
        ("DataLens", "исторические винрейты и H2H"),
        ("Magic Oculus", "архив турниров"),
        ("GitHub Actions", "debug/prod deployment"),
    ]
    for i, (h, b) in enumerate(integrations):
        y = 2.88 + i * 0.72
        txt(s, h, 6.78, y, 1.45, 0.24, 12, CYAN, True)
        txt(s, b, 8.38, y, 3.45, 0.25, 12, WHITE)
    txt(
        s,
        "Тестируется бизнес-логика, а не Telegram-моки",
        6.75,
        5.72,
        5.2,
        0.28,
        14,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 8)

    # 9 dirty data
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Самая инженерная часть", "Превращение грязных внешних данных в устойчивую модель")
    items = [
        ("Два HTML-формата", "активный и завершённый AetherHub выглядят по-разному", BLUE),
        ("Нормализация имён", "порядок имени, отчество, points внутри строки, bye", CYAN),
        ("Канонизация колод", "разные написания → общий архетип", GOLD),
        ("Контроль полноты", "roster, места и один head миграций проверяются", GREEN),
    ]
    for i, (h, b, c) in enumerate(items):
        x = 0.75 + (i % 2) * 6.05
        y = 2.05 + (i // 2) * 2.02
        box(s, x, y, 5.55, 1.58, PANEL, line=c)
        txt(s, h, x + 0.23, y + 0.29, 4.95, 0.3, 17, c, True)
        txt(s, b, x + 0.23, y + 0.83, 4.95, 0.45, 12, WHITE)
    txt(
        s,
        "Результат: система объяснимо деградирует и не превращает мусорный input в ложную статистику.",
        1.0,
        6.22,
        11.25,
        0.45,
        17,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 9)

    # 10 output
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Результат", "Закрытие турнира автоматически создаёт готовый медиапакет")
    add_picture_contain(s, ASSETS / "real-meta-chart.png", 0.7, 1.82, 3.85, 5.15, "0D1016")
    add_picture_contain(s, ASSETS / "real-standings.png", 4.74, 1.82, 3.85, 5.15, "0D1016")
    add_picture_contain(s, ASSETS / "real-telegram-recap.png", 8.78, 1.82, 3.85, 5.15, "F3F4F6")
    txt(s, "метагейм", 1.1, 6.64, 3.05, 0.22, 11, CYAN, True, align=PP_ALIGN.CENTER)
    txt(s, "финальные места", 5.12, 6.64, 3.05, 0.22, 11, GOLD, True, align=PP_ALIGN.CENTER)
    txt(s, "пост в чат", 9.15, 6.64, 3.05, 0.22, 11, GREEN, True, align=PP_ALIGN.CENTER)
    footer(s, 10, "Реальные материалы турнира на 42 участника")

    # 11 value flywheel
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Продуктовый эффект", "Данные создают петлю полезности, а не архив ради архива")
    flow = [
        ("Простая запись", "больше игроков\nуказывают колоду", BLUE),
        ("Полнее данные", "точнее локальная\nкартина меты", CYAN),
        ("Полезный ответ", "графики, H2H,\nистория игрока", GOLD),
        ("Больше доверия", "люди снова\nпользуются ботом", GREEN),
    ]
    for i, (h, b, c) in enumerate(flow):
        x = 0.75 + i * 3.1
        box(s, x, 2.25, 2.72, 2.15, PANEL, line=c)
        txt(s, h, x + 0.18, 2.62, 2.35, 0.34, 17, c, True, align=PP_ALIGN.CENTER)
        txt(s, b, x + 0.18, 3.28, 2.35, 0.65, 13, WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, "→", x + 2.75, 3.02, 0.32, 0.3, 18, GOLD, True)
    txt(
        s,
        "↖──────────────────────────────────────────────────────────────↙",
        1.15,
        4.72,
        11.0,
        0.32,
        15,
        "586273",
        align=PP_ALIGN.CENTER,
    )
    add_metric(s, 2.2, 5.38, "513 / 514", "реальных колод классифицированы", CYAN, 2.65)
    add_metric(s, 5.35, 5.38, "145", "архетипов в справочнике", GOLD, 2.65)
    add_metric(s, 8.5, 5.38, "42", "участника в показанном отчёте", BLUE, 2.65)
    footer(s, 11, "Метрики классификации: docs/meta_chart.md")

    # 12 safety
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Инженерная зрелость", "Безопасность здесь — часть предметной области")
    safeguards = [
        ("Никаких массовых debug-DM", "тестовый сценарий пишет только инициатору", RED),
        ("Opt-in уведомлений", "игрок явно включает полезные пуши", CYAN),
        ("Allowlist в production", "получатель проверяется перед отправкой", GOLD),
        ("Deploy через CI", "debug и prod разведены, секреты не в git", GREEN),
    ]
    for i, (h, b, c) in enumerate(safeguards):
        x = 0.78 + (i % 2) * 6.0
        y = 2.0 + (i // 2) * 2.1
        box(s, x, y, 5.5, 1.65, PANEL, line=c)
        txt(s, h, x + 0.25, y + 0.3, 5.0, 0.32, 17, c, True)
        txt(s, b, x + 0.25, y + 0.88, 4.95, 0.38, 12, WHITE)
    txt(
        s,
        "Ошибка fan-out спамит реальных людей и необратима — поэтому safety-инварианты закреплены в проектных правилах.",
        0.95,
        6.34,
        11.4,
        0.4,
        15,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 12)

    # 13 close
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Вывод", "Хороший community-tech невидим, пока не начинает помогать")
    txt(
        s,
        "MetaGatherer превращает разговор в чате\nв надёжное событие и полезные данные.",
        0.82,
        2.05,
        7.3,
        1.25,
        30,
        WHITE,
        True,
    )
    outcomes = [
        ("Игрок", "быстрее включается в турнир", CYAN),
        ("Организатор", "тратит меньше внимания на рутину", GOLD),
        ("Комьюнити", "получает общую память и аналитику", BLUE),
    ]
    for i, (h, b, c) in enumerate(outcomes):
        y = 3.72 + i * 0.85
        txt(s, h, 0.88, y, 1.45, 0.28, 14, c, True)
        txt(s, b, 2.45, y, 5.3, 0.3, 15, WHITE)
    s.shapes.add_picture(str(ASSETS / "real-final-post.jpg"), Inches(8.65), Inches(1.35), width=Inches(3.9))
    footer(s, 13)

    prs.core_properties.title = "MetaGatherer — подробная презентация"
    prs.core_properties.author = "MetaGatherer"
    prs.save(PPTX_OUT)


def video_frame(path: Path, eyebrow: str, headline: str, body: str, image_path: Path | None, accent: str):
    im = Image.new("RGB", (1920, 1080), "#10141C")
    d = ImageDraw.Draw(im)
    d.text((110, 78), eyebrow.upper(), font=font(26, True), fill=f"#{accent}")
    lines = headline.split("\n")
    y = 145
    for line in lines:
        d.text((110, y), line, font=font(68, True), fill="#F7F3E8")
        y += 82
    by = 885
    d.text((110, by), body, font=font(30), fill="#AAB3C2")
    if image_path:
        with Image.open(image_path).convert("RGB") as src:
            src.thumbnail((760, 720))
            x = 1080 + (760 - src.width) // 2
            iy = 165 + (720 - src.height) // 2
            d.rounded_rectangle(
                (x - 18, iy - 18, x + src.width + 18, iy + src.height + 18),
                28,
                fill="#202837",
                outline=f"#{accent}",
                width=4,
            )
            im.paste(src, (x, iy))
    else:
        # Abstract playing cards: enough context without using copyrighted card art.
        cards = [(1260, 260, -8, "#5591E6"), (1405, 220, 0, "#D6A928"), (1550, 270, 8, "#E4553D")]
        for x, cy, _, color in cards:
            d.rounded_rectangle((x, cy, x + 250, cy + 360), 24, fill="#202837", outline=color, width=8)
            d.ellipse((x + 72, cy + 70, x + 178, cy + 176), fill=color)
            d.rectangle((x + 45, cy + 230, x + 205, cy + 244), fill="#697486")
            d.rectangle((x + 45, cy + 270, x + 175, cy + 282), fill="#505A69")
        d.text((1370, 680), "MAGIC: THE GATHERING", font=font(24, True), fill="#AAB3C2")
    im.save(path, quality=94)


def build_video():
    VIDEO_DIR.mkdir(parents=True, exist_ok=True)
    scenes = [
        (
            "С чего всё начинается",
            "Есть карточная игра\nMagic: The Gathering",
            "Игроки собирают колоды и встречаются на клубных турнирах",
            None,
            GOLD,
        ),
        (
            "Разные стратегии",
            "Каждая колода играет\nпо-своему",
            "Tolarian Terror и Experimental Synthesizer — реальные карты из разных стратегий",
            CARDS_IMAGE,
            BLUE,
        ),
        (
            "Что такое метагейм",
            "Важно знать,\nкто чем играл",
            "Распределение колод на турнирах и есть локальный метагейм",
            ASSETS / "real-telegram-recap.png",
            CYAN,
        ),
        (
            "Как собираем данные",
            "Игрок указывает\nсвою колоду в Telegram",
            "Несколько нажатий — без формы и отдельного приложения",
            ASSETS / "real-choose-deck.jpg",
            BLUE,
        ),
        (
            "Если данных не хватает",
            "Организатор помогает\nзаполнить пробелы",
            "В итоге для каждого участника остаётся связка: игрок → турнир → колода",
            ASSETS / "real-admin-fill.jpg",
            GOLD,
        ),
        (
            "Весь процесс",
            "Бот связывает\nсемь шагов",
            "После публикации турнир обновляет общую мету, матчи и винрейты",
            FLOW_IMAGE,
            GOLD,
        ),
        (
            "История метагейма",
            "Становится видна\nпопулярность колод",
            "Сколько участий, турниров и какую долю локальной меты занимает стратегия",
            ASSETS / "real-meta-history.png",
            CYAN,
        ),
        (
            "Статистика результатов",
            "Можно считать\nматчи и винрейты",
            "Например: 77 матчей Blue Terror и 52% побед",
            ASSETS / "real-winrates.png",
            GREEN,
        ),
        (
            "Автоматика",
            "Пары и результаты\nприходят из AetherHub",
            "Повторный импорт обновляет данные без дублей",
            ASSETS / "Screenshot 2026-07-31 at 19.02.18.png",
            GREEN,
        ),
        (
            "Итог турнира",
            "Места и очки\nсобираются автоматически",
            "В одном изображении: игрок, колода, цветовая идентичность и результат",
            ASSETS / "real-standings.png",
            GOLD,
        ),
        (
            "После турнира",
            "Бот публикует\nготовый отчёт",
            "А данные турнира входят в общую статистику меты и результатов",
            ASSETS / "real-telegram-recap.png",
            RED,
        ),
        (
            "Зачем это людям",
            "Меньше рутины.\nБольше понимания игры.",
            "Игроки видят тренды и винрейты, а сообщество сохраняет общую историю",
            ASSETS / "real-meta-chart.png",
            GOLD,
        ),
    ]
    # The AetherHub scene uses a copied screenshot when available; fall back safely.
    fallback = ASSETS / "real-registration.jpg"
    frames = []
    for i, (e, h, b, img, c) in enumerate(scenes):
        if img and not img.exists():
            img = fallback
        p = VIDEO_DIR / f"scene-{i:02d}.jpg"
        if img == FLOW_IMAGE:
            with Image.open(FLOW_IMAGE).convert("RGB") as flow:
                flow.resize((1920, 1080), Image.Resampling.LANCZOS).save(p, quality=95)
        else:
            video_frame(p, e, h, b, img, c)
        frames.append(p)
    concat = VIDEO_DIR / "concat.txt"
    lines = []
    for p in frames:
        lines += [f"file '{p.resolve()}'", "duration 6"]
    lines.append(f"file '{frames[-1].resolve()}'")
    concat.write_text("\n".join(lines) + "\n", encoding="utf-8")
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(concat),
            "-vf",
            "fps=30,format=yuv420p,fade=t=in:st=0:d=0.5,fade=t=out:st=77:d=0.8",
            "-c:v",
            "libx264",
            "-preset",
            "medium",
            "-crf",
            "19",
            "-movflags",
            "+faststart",
            str(VIDEO_OUT),
        ],
        check=True,
    )


if __name__ == "__main__":
    build_supporting_images()
    build_flow_image()
    build_pptx()
    build_video()
    print(PPTX_OUT)
    print(VIDEO_OUT)
