"""Build a short Russian presentation about MetaGatherer.

Run:
    PYTHONPATH=/private/tmp/metagatherer-presentation python3 docs/presentation/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "MetaGatherer_short_ru.pptx"

W, H = 13.333, 7.5
BG = "10141C"
PANEL = "171D28"
PANEL_2 = "202837"
WHITE = "F7F3E8"
MUTED = "AAB3C2"
GOLD = "D6A928"
CYAN = "4FC3B3"
BLUE = "5591E6"
RED = "E4553D"
GREEN = "55C787"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def box(slide, x, y, w, h, fill=PANEL, radius=True, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def txt(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=WHITE,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def title(slide, kicker, headline, sub=None):
    txt(slide, kicker.upper(), 0.65, 0.35, 12.0, 0.3, 9, GOLD, True)
    txt(slide, headline, 0.65, 0.73, 12.0, 0.72, 29, WHITE, True)
    if sub:
        txt(slide, sub, 0.68, 1.48, 11.6, 0.52, 13, MUTED)


def footer(slide, n, source="MetaGatherer · 2026"):
    txt(slide, source, 0.68, 7.14, 10.5, 0.2, 8, "6F7888")
    txt(slide, f"0{n}", 12.05, 7.10, 0.6, 0.22, 8, GOLD, True, align=PP_ALIGN.RIGHT)


def font(size, bold=False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def make_phone(path: Path, screen: str):
    """Render a privacy-safe reconstruction based on actual bot strings/keyboards."""
    im = Image.new("RGB", (690, 1320), "#0E1621")
    d = ImageDraw.Draw(im)
    d.rounded_rectangle((8, 8, 682, 1312), 58, fill="#111822", outline="#303A48", width=7)
    d.rounded_rectangle((220, 20, 470, 58), 20, fill="#05070A")
    d.rectangle((20, 72, 670, 170), fill="#17212B")
    d.ellipse((46, 91, 108, 153), fill="#D6A928")
    d.text((126, 92), "MetaGatherer", font=font(31, True), fill="#FFFFFF")
    d.text((126, 130), "бот", font=font(20), fill="#7F91A4")

    def bubble(y, lines, buttons=()):
        line_h = 39
        bh = 44 + line_h * len(lines) + (61 * len(buttons) if buttons else 0)
        d.rounded_rectangle((44, y, 640, y + bh), 26, fill="#182533")
        ty = y + 24
        for line, strong, color in lines:
            d.text((70, ty), line, font=font(24, strong), fill=color)
            ty += line_h
        for label in buttons:
            d.rounded_rectangle((70, ty + 8, 614, ty + 56), 14, fill="#2B5278")
            tw = d.textbbox((0, 0), label, font=font(21, True))[2]
            d.text((342 - tw / 2, ty + 20), label, font=font(21, True), fill="#FFFFFF")
            ty += 61

    if screen == "join":
        bubble(
            220,
            [
                ("🏆 Goldfish Pauper · сегодня", True, "#FFFFFF"),
                ("Регистрация · 18 человек", False, "#AFC1D3"),
                ("✅ 15 с колодой   ⬜ 3 без", False, "#AFC1D3"),
            ],
            ("🃏 Выбрать колоду", "🚪 Выйти из турнира"),
        )
        bubble(
            585, [("Выберите архетип колоды:", True, "#FFFFFF")], ("🔵 Blue Terror", "🔴 Red Madness", "… ещё колоды")
        )
    elif screen == "round":
        bubble(
            230,
            [
                ("⚔️ Раунд 3", True, "#FFFFFF"),
                ("Ваш оппонент: Алексей", False, "#FFFFFF"),
                ("Стол 4", False, "#AFC1D3"),
                ("Последние колоды: Affinity, Gates", False, "#78B7E8"),
                ("Ваш H2H: 58% · 12 партий", False, "#78B7E8"),
            ],
            ("👥 Показать оппонентов",),
        )
        bubble(
            650,
            [
                ("Данные обновлены из AetherHub", True, "#FFFFFF"),
                ("Следующая проверка — автоматически", False, "#AFC1D3"),
            ],
        )
    else:
        bubble(
            220,
            [
                ("🎉 Сбор метагейма завершён", True, "#FFFFFF"),
                ("Участников: 21 · колод собрано: 21", False, "#AFC1D3"),
                ("Без поражений: 2 игрока", False, "#AFC1D3"),
                ("Спасибо метаписцам 🙏", False, "#D6A928"),
            ],
            ("🍩 График метагейма", "📊 Итоговые стендинги"),
        )
        bubble(
            600,
            [
                ("Экспорт готов", True, "#FFFFFF"),
                ("CSV / Excel / Markdown", False, "#AFC1D3"),
            ],
        )
    im.save(path, quality=94)


def make_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    for name in ("join", "round", "finish"):
        make_phone(ASSETS / f"telegram-{name}.png", name)


def add_metric(slide, x, y, value, label, accent=GOLD, w=2.35):
    box(slide, x, y, w, 1.17, PANEL)
    txt(slide, value, x + 0.18, y + 0.15, w - 0.35, 0.45, 25, accent, True)
    txt(slide, label, x + 0.18, y + 0.67, w - 0.35, 0.3, 10, MUTED)


def main():
    make_assets()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 — promise
    s = prs.slides.add_slide(blank)
    add_bg(s)
    s.shapes.add_picture(
        str(ASSETS / "bot-avatar.png"), Inches(0.72), Inches(0.34), width=Inches(0.62), height=Inches(0.62)
    )
    txt(s, "METAGATHERER", 1.48, 0.52, 4.0, 0.3, 10, GOLD, True)
    txt(s, "Турнир живёт\nпрямо в Telegram", 0.72, 1.12, 7.2, 1.55, 35, WHITE, True)
    txt(
        s,
        "Регистрация, колоды, пары, результаты и мета —\nбез таблиц, ручных напоминаний и потери данных.",
        0.76,
        2.92,
        6.4,
        0.9,
        16,
        MUTED,
    )
    box(s, 0.75, 4.35, 5.85, 1.18, PANEL)
    txt(s, "Для людей", 1.0, 4.62, 1.2, 0.28, 11, CYAN, True)
    txt(s, "меньше рутины · больше игры", 1.0, 4.94, 4.7, 0.3, 16, WHITE, True)
    box(s, 0.75, 5.72, 5.85, 1.18, PANEL)
    txt(s, "Для комьюнити", 1.0, 5.99, 1.65, 0.28, 11, GOLD, True)
    txt(s, "достоверная история локальной меты", 1.0, 6.31, 4.9, 0.3, 16, WHITE, True)
    s.shapes.add_picture(str(ASSETS / "real-telegram-recap.png"), Inches(8.48), Inches(0.62), height=Inches(6.58))
    footer(s, 1)

    # 2 — problem
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Зачем", "Локальный турнир — это маленькая распределённая система")
    problems = [
        ("01", "Люди", "приходят через чат,\nчасто в последний момент", BLUE),
        ("02", "Колоды", "записывают игроки,\nадмины и оппоненты", CYAN),
        ("03", "Результаты", "живут во внешнем\nтурнирном сервисе", GOLD),
        ("04", "История", "теряется между чатами,\nтаблицами и дашбордом", RED),
    ]
    for i, (num, head, body, c) in enumerate(problems):
        x = 0.68 + i * 3.12
        box(s, x, 2.25, 2.83, 2.45, PANEL)
        txt(s, num, x + 0.2, 2.45, 0.55, 0.35, 14, c, True)
        txt(s, head, x + 0.2, 2.96, 2.4, 0.36, 19, WHITE, True)
        txt(s, body, x + 0.2, 3.53, 2.35, 0.78, 13, MUTED)
    txt(
        s,
        "MetaGatherer связывает всё одним идентификатором турнира и понятным сценарием в чате.",
        0.76,
        5.34,
        11.7,
        0.55,
        21,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    txt(
        s,
        "Telegram становится интерфейсом · база — памятью · автоматика — оператором",
        1.4,
        6.12,
        10.5,
        0.34,
        13,
        GOLD,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 2)

    # 3 — user journey
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(
        s,
        "Как это выглядит",
        "От выбора колоды до готового отчёта — в одном чате",
        "Реальные кадры из рабочего Telegram-сценария.",
    )
    labels = [
        ("1", "Игрок", "выбирает колоду"),
        ("2", "Организатор", "закрывает пробелы"),
        ("3", "Комьюнити", "получает итог"),
    ]
    real_images = ("real-choose-deck.jpg", "real-admin-fill.jpg", "real-final-post.jpg")
    for i, (n, head, body) in enumerate(labels):
        x = 0.66 + i * 4.23
        box(s, x, 2.10, 3.72, 3.55, "E9ECF1", line="313A48")
        s.shapes.add_picture(str(ASSETS / real_images[i]), Inches(x + 0.08), Inches(2.18), width=Inches(3.56))
        box(s, x, 5.88, 3.72, 0.74, PANEL)
        txt(s, n, x + 0.15, 6.09, 0.34, 0.22, 11, GOLD, True, align=PP_ALIGN.CENTER)
        txt(s, head, x + 0.55, 6.02, 1.52, 0.25, 13, WHITE, True)
        txt(s, body, x + 1.95, 6.06, 1.55, 0.22, 9, MUTED, align=PP_ALIGN.RIGHT)
    footer(s, 3, "Реальные кадры Telegram Desktop · июль 2026")

    # 4 — architecture
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Технически круто", "Не просто бот: устойчивый data pipeline вокруг события")
    nodes = [
        (0.65, "Telegram", "тонкие async-обёртки", BLUE),
        (3.15, "Handlers", "чистая бизнес-логика", CYAN),
        (5.65, "Services", "state machine · правила", GOLD),
        (8.15, "PostgreSQL", "единая история", GREEN),
        (10.65, "Визуализация", "PNG · Excel · Web", RED),
    ]
    for i, (x, head, body, c) in enumerate(nodes):
        box(s, x, 2.26, 2.02, 1.38, PANEL, line=c)
        txt(s, head, x + 0.15, 2.55, 1.72, 0.28, 15, WHITE, True, align=PP_ALIGN.CENTER)
        txt(s, body, x + 0.14, 3.02, 1.75, 0.24, 9, MUTED, align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            txt(s, "→", x + 2.08, 2.72, 0.35, 0.3, 18, "6F7888", True, align=PP_ALIGN.CENTER)
    cards = [
        (0.8, "Деградация без падения", "AetherHub/DataLens недоступны — основной сценарий продолжает работать."),
        (4.38, "Безопасные уведомления", "opt-in, allowlist и запрет массовых debug-DM защищают реальных людей."),
        (7.96, "Грязные данные → модель", "имена, колоды и два HTML-формата нормализуются и проверяются."),
    ]
    for x, head, body in cards:
        box(s, x, 4.43, 3.18, 1.72, PANEL_2)
        txt(s, head, x + 0.22, 4.69, 2.75, 0.33, 14, GOLD, True)
        txt(s, body, x + 0.22, 5.18, 2.72, 0.7, 11, MUTED)
    txt(
        s,
        "Интеграции: AetherHub · DataLens · Magic Oculus · GitHub Actions",
        2.1,
        6.55,
        9.1,
        0.3,
        12,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
    )
    footer(s, 4, "Архитектура и инварианты — по коду и docs/")

    # 5 — data
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Данные становятся полезными", "Каждый дейлик оставляет цифровой след")
    # Use the real generated chart supplied with the product demo.
    ref = Image.open(ASSETS / "real-meta-chart.png")
    crop = ref.crop((0, 0, ref.width, int(ref.height * 0.72)))
    crop_path = ASSETS / "meta-crop.png"
    crop.save(crop_path)
    box(s, 0.68, 1.83, 4.35, 4.95, "0D1016")
    s.shapes.add_picture(str(crop_path), Inches(1.04), Inches(1.98), width=Inches(3.62), height=Inches(4.65))
    txt(s, "Реальный отчёт: 42 колоды", 1.22, 6.55, 3.2, 0.2, 8, "788191", align=PP_ALIGN.CENTER)
    txt(s, "Сбор", 5.55, 2.03, 1.3, 0.28, 12, CYAN, True)
    txt(s, "игрок → колода → пары → место", 5.55, 2.38, 6.4, 0.45, 20, WHITE, True)
    txt(s, "↓", 5.56, 2.96, 0.5, 0.34, 20, GOLD, True)
    txt(s, "Нормализация", 5.55, 3.36, 1.8, 0.28, 12, CYAN, True)
    txt(s, "разные написания сводятся к общему архетипу", 5.55, 3.71, 6.5, 0.45, 18, WHITE, True)
    txt(s, "↓", 5.56, 4.28, 0.5, 0.34, 20, GOLD, True)
    txt(s, "Обратная связь", 5.55, 4.68, 2.0, 0.28, 12, CYAN, True)
    txt(s, "мета · стендинги · H2H · история игрока", 5.55, 5.03, 6.5, 0.45, 18, WHITE, True)
    add_metric(s, 5.55, 5.83, "513 / 514", "реальных колод получили цвет", CYAN, 2.35)
    add_metric(s, 8.08, 5.83, "145", "архетипов в справочнике", GOLD, 2.0)
    add_metric(s, 10.26, 5.83, "3", "источника данных", BLUE, 2.0)
    footer(s, 5, "Реальный график из Telegram · метрики: docs/meta_chart.md")

    # 6 — outcome
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Итог", "Автоматизация, которую сообщество действительно замечает")
    outcomes = [
        ("Игроку", "записаться за минуту,\nне пропустить пару,\nувидеть свой контекст", CYAN),
        ("Организатору", "не собирать колоды вручную,\nне сводить таблицы,\nбыстро закрыть турнир", GOLD),
        ("Сообществу", "видеть живую мету,\nсохранять историю,\nпринимать решения по данным", BLUE),
    ]
    for i, (head, body, c) in enumerate(outcomes):
        x = 0.72 + i * 4.18
        box(s, x, 2.08, 3.72, 2.58, PANEL, line=c)
        txt(s, head, x + 0.25, 2.42, 3.2, 0.4, 21, c, True)
        txt(s, body, x + 0.25, 3.12, 3.05, 1.15, 15, WHITE)
    txt(s, "Главная идея", 0.78, 5.26, 1.45, 0.25, 11, GOLD, True)
    txt(
        s,
        "MetaGatherer превращает разговор в Telegram в надёжные данные —\nи возвращает эти данные людям в момент, когда они полезны.",
        0.78,
        5.66,
        11.5,
        0.92,
        24,
        WHITE,
        True,
    )
    footer(s, 6)

    prs.core_properties.title = "MetaGatherer — коротко о проекте"
    prs.core_properties.subject = "Презентация для технической аудитории, не знакомой с MTG"
    prs.core_properties.author = "MetaGatherer"
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
